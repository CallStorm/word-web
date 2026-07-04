"""上传文档 → 纯文本/markdown 提取（在 backend 进程内，不写盘）。

用于「智能填充」文档模式：把用户上传的素材解析成文本，喂给 LLM 一次性生成
主题 / 大纲 / 重点 / 推荐设置 / 设计风格。

复用 ppt-master source_to_md 同款库（PyMuPDF / mammoth / python-pptx / openpyxl /
markdownify / beautifulsoup4），但只取文本、丢弃图片与表格渲染细节，避免写盘。

不支持图片/扫描件（无 OCR）：遇到这类文件抛 DocumentExtractError，由上层转成
对用户友好的提示。
"""
from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path

log = logging.getLogger("backend.app.document_extract")


class DocumentExtractError(Exception):
    """文档解析失败（如格式不支持、图片无 OCR、解析异常）。"""


# 单文件文本上限与整体预算，避免超长文档把 LLM 上下文撑爆
_PER_FILE_CHAR_LIMIT = 20_000
_TOTAL_CHAR_LIMIT = 30_000


def _truncate(text: str, limit: int) -> str:
    text = text.strip()
    if len(text) <= limit:
        return text
    return text[:limit] + "\n…（已截断）"


def _extract_pdf(content: bytes) -> str:
    import fitz  # PyMuPDF

    out: list[str] = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        for page in doc:
            out.append(page.get_text("text"))
    return "\n\n".join(out)


def _extract_docx(content: bytes) -> str:
    import mammoth

    # mammoth 输出 markdown，保留标题层级
    result = mammoth.convert_to_markdown(BytesIO(content))
    return result.value


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    out: list[str] = []
    pres = Presentation(BytesIO(content))
    for idx, slide in enumerate(pres.slides, start=1):
        chunks: list[str] = [f"## 第 {idx} 页"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = "\n".join(
                p.text for p in shape.text_frame.paragraphs if p.text
            )
            if txt.strip():
                chunks.append(txt)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                chunks.append(f"（备注：{notes}）")
        out.append("\n".join(chunks))
    return "\n\n".join(out)


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    out: list[str] = []
    wb = load_workbook(BytesIO(content), data_only=True, read_only=True)
    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            out.append(f"## {ws.title}\n" + "\n".join(rows[:200]))
    wb.close()
    return "\n\n".join(out)


def _extract_html(content: bytes) -> str:
    from bs4 import BeautifulSoup
    from markdownify import markdownify as md

    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return md(str(soup), heading_style="ATX")


def _extract_text(content: bytes) -> str:
    # md / txt / csv 等：直接按 utf-8 解码（失败再尝试常见编码）
    for enc in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return ""


# 扩展名 → 提取器
_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".md": _extract_text,
    ".markdown": _extract_text,
    ".txt": _extract_text,
    ".csv": _extract_text,
}


def extract_one(name: str, content: bytes) -> str:
    """解析单个文件 → markdown 文本。失败抛 DocumentExtractError。"""
    suffix = Path(name).suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        # 图片 / epub / 其它未接入格式：明确报错，交给上层提示
        raise DocumentExtractError(
            f"暂不支持从 {suffix or '该'} 文件提取文本"
            "（图片/扫描件等暂不支持，请改用文本类文档，或在主题模式下手动填写）"
        )
    try:
        text = extractor(content)
    except DocumentExtractError:
        raise
    except Exception as e:  # 解析库自身的异常
        log.warning("解析 %s 失败: %s", name, e)
        raise DocumentExtractError(f"解析 {name} 失败：{e}") from e
    text = _truncate(text, _PER_FILE_CHAR_LIMIT)
    if not text:
        raise DocumentExtractError(
            f"{name} 未提取到任何文本（可能是图片/扫描件，暂不支持）"
        )
    return text


def extract_document_text(files: list[tuple[str, bytes]]) -> str:
    """解析多个上传文件，拼接成一段 markdown，截断到整体预算。

    files: [(filename, raw_bytes), ...]
    返回带文件分隔的 markdown 文本；任一文件解析失败即抛 DocumentExtractError
    （整体失败，不部分填充）。
    """
    if not files:
        raise DocumentExtractError("未上传任何文件")

    parts: list[str] = []
    total = 0
    for name, content in files:
        md = extract_one(name, content)
        chunk = f"### 文件：{name}\n{md}"
        # 整体预算控制：超了就停止追加后续文件
        if total + len(chunk) > _TOTAL_CHAR_LIMIT:
            remaining = _TOTAL_CHAR_LIMIT - total
            if remaining > 200:
                parts.append(chunk[:remaining] + "\n…（已截断）")
            parts.append("\n…（文档过长，部分文件已省略）")
            break
        parts.append(chunk)
        total += len(chunk)

    result = "\n\n".join(parts).strip()
    if not result:
        raise DocumentExtractError("所有文件均未提取到文本")
    return result
